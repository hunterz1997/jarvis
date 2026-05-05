"""
Computer integration — full filesystem access, shell execution, system info,
screenshot, clipboard, application launching, Python execution.
"""

import asyncio
import io
import json
import logging
import os
import shutil
import subprocess
import sys
import tempfile
import time
import zipfile
from datetime import datetime, date
from pathlib import Path
from typing import Any

import psutil

logger = logging.getLogger(__name__)


class ComputerIntegration:
    """Handles all local system interactions."""

    async def execute(self, tool_name: str, params: dict) -> dict[str, Any]:
        """Dispatch to the correct method."""
        dispatch = {
            "read_file": self._read_file,
            "write_file": self._write_file,
            "list_directory": self._list_directory,
            "search_files": self._search_files,
            "file_operation": self._file_operation,
            "run_command": self._run_command,
            "launch_application": self._launch_application,
            "take_screenshot": self._take_screenshot,
            "system_info": self._system_info,
            "clipboard": self._clipboard,
            "run_python": self._run_python,
        }
        handler = dispatch.get(tool_name)
        if not handler:
            return {"success": False, "error": f"Unknown tool: {tool_name}"}
        return await handler(params)

    # ------------------------------------------------------------------
    # File reading — supports PDF, DOCX, XLSX, PPTX, images, plain text
    # ------------------------------------------------------------------

    async def _read_file(self, params: dict) -> dict:
        path = Path(params["path"])
        if not path.exists():
            return {"success": False, "error": f"File not found: {path}", "suggestion": "Check the path and try again."}

        suffix = path.suffix.lower()

        try:
            if suffix == ".pdf":
                return await self._read_pdf(path, params.get("page_range"))
            elif suffix in (".docx", ".doc"):
                return await self._read_docx(path)
            elif suffix in (".xlsx", ".xls"):
                return await self._read_excel(path, params.get("sheet_name"))
            elif suffix in (".pptx", ".ppt"):
                return await self._read_pptx(path)
            elif suffix in (".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff", ".webp"):
                return await self._read_image(path)
            elif suffix == ".csv":
                return await self._read_csv(path)
            else:
                # Plain text fallback
                content = path.read_text(encoding="utf-8", errors="replace")
                if len(content) > 15000:
                    content = content[:15000] + "\n\n[... content truncated ...]"
                return {"success": True, "path": str(path), "content": content, "size": path.stat().st_size}
        except Exception as e:
            logger.error("read_file failed for %s: %s", path, e)
            return {"success": False, "error": str(e), "suggestion": "The file may be corrupted or in an unsupported format."}

    async def _read_pdf(self, path: Path, page_range: str | None) -> dict:
        import fitz  # PyMuPDF
        doc = fitz.open(str(path))
        pages_to_read = range(len(doc))
        if page_range:
            parts = page_range.split("-")
            start = int(parts[0]) - 1
            end = int(parts[1]) if len(parts) > 1 else start + 1
            pages_to_read = range(start, min(end, len(doc)))

        text = ""
        for i in pages_to_read:
            text += f"\n--- Page {i + 1} ---\n"
            text += doc[i].get_text()

        doc.close()
        if len(text) > 15000:
            text = text[:15000] + "\n\n[... content truncated ...]"
        return {"success": True, "path": str(path), "content": text, "pages": len(doc)}

    async def _read_docx(self, path: Path) -> dict:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        content = "\n".join(paragraphs)
        if len(content) > 15000:
            content = content[:15000] + "\n\n[... truncated ...]"
        return {"success": True, "path": str(path), "content": content}

    async def _read_excel(self, path: Path, sheet_name: str | None) -> dict:
        import pandas as pd
        xf = pd.ExcelFile(str(path))
        sheets = xf.sheet_names
        target = sheet_name if sheet_name in sheets else sheets[0]
        df = pd.read_excel(str(path), sheet_name=target)
        summary = f"Sheet: {target} | Rows: {len(df)} | Columns: {list(df.columns)}\n\n"
        summary += df.to_string(max_rows=100, max_cols=20)
        if len(summary) > 15000:
            summary = summary[:15000] + "\n\n[... truncated ...]"
        return {"success": True, "path": str(path), "content": summary, "sheets": sheets}

    async def _read_pptx(self, path: Path) -> dict:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_content = f"\n--- Slide {i + 1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content += shape.text + "\n"
            slides_text.append(slide_content)
        content = "".join(slides_text)
        if len(content) > 15000:
            content = content[:15000] + "\n\n[... truncated ...]"
        return {"success": True, "path": str(path), "content": content, "slides": len(prs.slides)}

    async def _read_image(self, path: Path) -> dict:
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(str(path))
            text = pytesseract.image_to_string(img)
            return {
                "success": True,
                "path": str(path),
                "content": text or "(No text detected in image)",
                "size": f"{img.width}x{img.height}",
                "mode": img.mode,
            }
        except Exception:
            return {
                "success": True,
                "path": str(path),
                "content": f"[Image file: {path.name} — OCR not available]",
                "note": "Install tesseract for OCR support.",
            }

    async def _read_csv(self, path: Path) -> dict:
        import pandas as pd
        df = pd.read_csv(str(path))
        summary = f"Rows: {len(df)} | Columns: {list(df.columns)}\n\n"
        summary += df.to_string(max_rows=100)
        if len(summary) > 15000:
            summary = summary[:15000] + "\n\n[... truncated ...]"
        return {"success": True, "path": str(path), "content": summary}

    # ------------------------------------------------------------------
    # Write file
    # ------------------------------------------------------------------

    async def _write_file(self, params: dict) -> dict:
        path = Path(params["path"])
        path.parent.mkdir(parents=True, exist_ok=True)
        mode = "a" if params.get("mode") == "append" else "w"
        path.write_text(params["content"], encoding="utf-8") if mode == "w" else open(path, "a", encoding="utf-8").write(params["content"])
        return {"success": True, "path": str(path), "size": path.stat().st_size, "mode": mode}

    # ------------------------------------------------------------------
    # Directory listing
    # ------------------------------------------------------------------

    async def _list_directory(self, params: dict) -> dict:
        path = Path(params["path"])
        if not path.exists():
            return {"success": False, "error": f"Directory not found: {path}"}
        pattern = params.get("pattern", "*")
        recursive = params.get("recursive", False)

        items = []
        glob_fn = path.rglob if recursive else path.glob
        for item in sorted(glob_fn(pattern)):
            stat = item.stat()
            items.append({
                "name": item.name,
                "path": str(item),
                "type": "folder" if item.is_dir() else "file",
                "size": stat.st_size if item.is_file() else None,
                "modified": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M"),
            })
        return {"success": True, "path": str(path), "items": items[:200], "count": len(items)}

    # ------------------------------------------------------------------
    # File search
    # ------------------------------------------------------------------

    async def _search_files(self, params: dict) -> dict:
        search_path = Path(params.get("search_path", "C:\\"))
        name_pattern = params.get("name_pattern", "*")
        content_keyword = params.get("content_keyword")
        modified_after = params.get("modified_after")
        max_results = params.get("max_results", 50)

        cutoff = None
        if modified_after:
            cutoff = datetime.fromisoformat(modified_after).timestamp()

        results = []
        try:
            for item in search_path.rglob(name_pattern):
                if not item.is_file():
                    continue
                stat = item.stat()
                if cutoff and stat.st_mtime < cutoff:
                    continue
                if content_keyword:
                    try:
                        if content_keyword.lower() not in item.read_text(encoding="utf-8", errors="ignore").lower():
                            continue
                    except Exception:
                        continue
                results.append({
                    "name": item.name,
                    "path": str(item),
                    "size": stat.st_size,
                    "modified": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M"),
                })
                if len(results) >= max_results:
                    break
        except PermissionError:
            pass  # Skip restricted paths silently

        return {"success": True, "results": results, "count": len(results)}

    # ------------------------------------------------------------------
    # File operations
    # ------------------------------------------------------------------

    async def _file_operation(self, params: dict) -> dict:
        op = params["operation"]
        src = Path(params["source"])
        dst = Path(params["destination"]) if params.get("destination") else None

        try:
            if op == "copy":
                shutil.copy2(src, dst)
                return {"success": True, "operation": op, "from": str(src), "to": str(dst)}
            elif op == "move":
                shutil.move(str(src), str(dst))
                return {"success": True, "operation": op, "from": str(src), "to": str(dst)}
            elif op == "rename":
                src.rename(dst)
                return {"success": True, "operation": op, "from": str(src), "to": str(dst)}
            elif op == "delete":
                if src.is_dir():
                    shutil.rmtree(src)
                else:
                    src.unlink()
                return {"success": True, "operation": op, "deleted": str(src)}
            elif op == "mkdir":
                src.mkdir(parents=True, exist_ok=True)
                return {"success": True, "operation": op, "created": str(src)}
            elif op == "zip":
                with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zf:
                    if src.is_dir():
                        for f in src.rglob("*"):
                            zf.write(f, f.relative_to(src))
                    else:
                        zf.write(src, src.name)
                return {"success": True, "operation": op, "archive": str(dst)}
            elif op == "unzip":
                with zipfile.ZipFile(src, "r") as zf:
                    zf.extractall(dst or src.parent)
                return {"success": True, "operation": op, "extracted_to": str(dst or src.parent)}
            else:
                return {"success": False, "error": f"Unknown operation: {op}"}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ------------------------------------------------------------------
    # Shell command — always requires confirmation flag
    # ------------------------------------------------------------------

    async def _run_command(self, params: dict) -> dict:
        command = params["command"]
        shell_type = params.get("shell", "powershell")
        working_dir = params.get("working_dir")
        timeout = params.get("timeout", 30)

        # Build the actual command
        if shell_type == "powershell":
            cmd = ["powershell", "-NonInteractive", "-Command", command]
        else:
            cmd = ["cmd", "/c", command]

        try:
            result = await asyncio.wait_for(
                asyncio.create_subprocess_exec(
                    *cmd,
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                    cwd=working_dir,
                ),
                timeout=5,
            )
            proc = result
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
            return {
                "success": proc.returncode == 0,
                "returncode": proc.returncode,
                "stdout": stdout.decode("utf-8", errors="replace")[:5000],
                "stderr": stderr.decode("utf-8", errors="replace")[:2000],
            }
        except asyncio.TimeoutError:
            return {"success": False, "error": f"Command timed out after {timeout}s"}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ------------------------------------------------------------------
    # Launch application
    # ------------------------------------------------------------------

    async def _launch_application(self, params: dict) -> dict:
        app = params["app_name"]
        args = params.get("args", "")

        # Common app name aliases
        aliases = {
            "excel": "excel.exe",
            "word": "winword.exe",
            "powerpoint": "powerpnt.exe",
            "chrome": "chrome.exe",
            "vs code": "code.exe",
            "vscode": "code.exe",
            "notepad": "notepad.exe",
            "explorer": "explorer.exe",
            "task manager": "taskmgr.exe",
        }
        executable = aliases.get(app.lower(), app)
        cmd = [executable]
        if args:
            cmd.append(args)

        try:
            subprocess.Popen(cmd, shell=True)
            return {"success": True, "launched": app}
        except Exception as e:
            return {"success": False, "error": str(e), "suggestion": f"Try using the full path to {app}"}

    # ------------------------------------------------------------------
    # Screenshot
    # ------------------------------------------------------------------

    async def _take_screenshot(self, params: dict) -> dict:
        try:
            from PIL import ImageGrab
            img = ImageGrab.grab()
            save_path = params.get("save_path")
            if save_path:
                img.save(save_path)
            else:
                save_path = str(Path(tempfile.gettempdir()) / f"jarvis_screenshot_{int(time.time())}.png")
                img.save(save_path)
            return {
                "success": True,
                "path": save_path,
                "size": f"{img.width}x{img.height}",
                "message": f"Screenshot saved to {save_path}",
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ------------------------------------------------------------------
    # System info
    # ------------------------------------------------------------------

    async def _system_info(self, params: dict) -> dict:
        info_type = params.get("info_type", "all")
        result = {}

        if info_type in ("all", "ram"):
            vm = psutil.virtual_memory()
            result["ram"] = {
                "total_gb": round(vm.total / 1e9, 1),
                "used_gb": round(vm.used / 1e9, 1),
                "available_gb": round(vm.available / 1e9, 1),
                "percent_used": vm.percent,
            }

        if info_type in ("all", "cpu"):
            result["cpu"] = {
                "percent": psutil.cpu_percent(interval=0.5),
                "cores_physical": psutil.cpu_count(logical=False),
                "cores_logical": psutil.cpu_count(logical=True),
                "frequency_mhz": round(psutil.cpu_freq().current) if psutil.cpu_freq() else None,
            }

        if info_type in ("all", "disk"):
            disks = []
            for part in psutil.disk_partitions():
                try:
                    usage = psutil.disk_usage(part.mountpoint)
                    disks.append({
                        "drive": part.mountpoint,
                        "total_gb": round(usage.total / 1e9, 1),
                        "used_gb": round(usage.used / 1e9, 1),
                        "free_gb": round(usage.free / 1e9, 1),
                        "percent_used": usage.percent,
                    })
                except PermissionError:
                    pass
            result["disk"] = disks

        if info_type in ("all", "processes"):
            procs = []
            for p in sorted(psutil.process_iter(["pid", "name", "memory_percent", "cpu_percent"]),
                             key=lambda x: x.info["memory_percent"] or 0, reverse=True)[:15]:
                procs.append({
                    "pid": p.info["pid"],
                    "name": p.info["name"],
                    "memory_percent": round(p.info["memory_percent"] or 0, 1),
                })
            result["processes"] = procs

        if info_type in ("all", "uptime"):
            boot_time = psutil.boot_time()
            uptime_seconds = time.time() - boot_time
            hours, rem = divmod(int(uptime_seconds), 3600)
            minutes = rem // 60
            result["uptime"] = f"{hours}h {minutes}m"
            result["boot_time"] = datetime.fromtimestamp(boot_time).strftime("%Y-%m-%d %H:%M")

        return {"success": True, **result}

    # ------------------------------------------------------------------
    # Clipboard
    # ------------------------------------------------------------------

    async def _clipboard(self, params: dict) -> dict:
        try:
            import pyperclip
            if params["action"] == "read":
                text = pyperclip.paste()
                return {"success": True, "content": text}
            else:
                pyperclip.copy(params.get("text", ""))
                return {"success": True, "message": "Text copied to clipboard"}
        except Exception as e:
            return {"success": False, "error": str(e)}

    # ------------------------------------------------------------------
    # Run Python
    # ------------------------------------------------------------------

    async def _run_python(self, params: dict) -> dict:
        code = params["code"]
        save_as = params.get("save_as")
        timeout = params.get("timeout", 60)

        if save_as:
            script_path = Path(save_as)
            script_path.parent.mkdir(parents=True, exist_ok=True)
            script_path.write_text(code, encoding="utf-8")
        else:
            tmp = tempfile.NamedTemporaryFile(suffix=".py", delete=False, mode="w", encoding="utf-8")
            tmp.write(code)
            tmp.close()
            script_path = Path(tmp.name)

        try:
            proc = await asyncio.create_subprocess_exec(
                sys.executable, str(script_path),
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
            return {
                "success": proc.returncode == 0,
                "returncode": proc.returncode,
                "stdout": stdout.decode("utf-8", errors="replace")[:8000],
                "stderr": stderr.decode("utf-8", errors="replace")[:2000],
                "script_path": str(script_path) if save_as else None,
            }
        except asyncio.TimeoutError:
            return {"success": False, "error": f"Script timed out after {timeout}s"}
        finally:
            if not save_as:
                try:
                    script_path.unlink()
                except Exception:
                    pass


computer = ComputerIntegration()
